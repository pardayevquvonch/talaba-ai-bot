import json
import sqlite3
import urllib.parse
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
import telebot
from telebot import types

# ==================== SOZLAMALAR ====================
BOT_TOKEN = "8935762970:AAHOJcuRmBeNd2Up4NdLVdPgjf8dzZgyolc"
GEMINI_API_KEY = "AIzaSy...Ab8RN6LepRLXF_p2-ufzFKvXKADqiib8MgVhQtsnw-4H7YZoig"  

genai.configure(api_key=GEMINI_API_KEY)
model = genai.GenerativeModel("gemini-1.5-flash")
bot = telebot.TeleBot(BOT_TOKEN)


# ==================== MA'LUMOTLAR BAZASI ====================
def init_db():
  conn = sqlite3.connect("users.db")
  cursor = conn.cursor()
  cursor.execute("""
    CREATE TABLE IF NOT EXISTS users (
        user_id INTEGER PRIMARY KEY,
        limits INTEGER DEFAULT 5
    )
    """)
  conn.commit()
  conn.close()


def get_user_limit(user_id):
  conn = sqlite3.connect("users.db")
  cursor = conn.cursor()
  cursor.execute("SELECT limits FROM users WHERE user_id = ?", (user_id,))
  row = cursor.fetchone()
  if row is None:
    cursor.execute(
        "INSERT INTO users (user_id, limits) VALUES (?, 5)", (user_id,)
    )
    conn.commit()
    conn.close()
    return 5
  conn.close()
  return row[0]


def update_user_limit(user_id, change):
  conn = sqlite3.connect("users.db")
  cursor = conn.cursor()
  cursor.execute(
      "UPDATE users SET limits = ? WHERE user_id = ?",
      (max(0, get_user_limit(user_id) + change), user_id),
  )
  conn.commit()
  conn.close()


init_db()
user_states = {}


# ==================== MENYU ====================
def main_menu():
  markup = types.ReplyKeyboardMarkup(row_width=2, resize_keyboard=True)
  btn1 = types.KeyboardButton("📊 AI Slayd Yaratish")
  btn2 = types.KeyboardButton("📝 AI Referat Yozish")
  btn3 = types.KeyboardButton("📚 Kurs Ishi Yozish")
  btn4 = types.KeyboardButton("📑 Mustaqil Ish")
  btn5 = types.KeyboardButton("✍️ Qo'lda Yozilgan Matn")
  btn6 = types.KeyboardButton("❓ Test / Quiz / Savol-Javob")
  btn7 = types.KeyboardButton("🌐 Til O'rganish (Language)")
  btn8 = types.KeyboardButton("📰 Maqola Yozish")
  btn9 = types.KeyboardButton("🎨 AI Rasm Yaratish")
  btn10 = types.KeyboardButton("⚡️ Mening Limitlarim")
  markup.add(btn1, btn2, btn3, btn4, btn5, btn6, btn7, btn8, btn9, btn10)
  return markup


# ==================== HANDLERLAR ====================
@bot.message_handler(commands=["start"])
def send_welcome(message):
  limit = get_user_limit(message.from_user.id)
  welcome_text = (
      f"Assalomu alaykum, <b>{message.from_user.first_name}</b>!\n\n"
      "🤖 <b>Bu bot nimalar qila oladi?</b>\n\n"
      "• 📊 <b>AI Slayd Yaratish</b> - Avtomatik PowerPoint (.pptx) taqdimot\n"
      "• 📝 <b>AI Referat Yozish</b> - Reja va xulosali tayyor matn\n"
      "• 📚 <b>Kurs Ishi Yozish</b> - Akademik va tuzilmali kurs ishi\n"
      "• 📑 <b>Mustaqil Ish</b> - Talabalar uchun mustaqil ish topshiriqlari\n"
      "• ✍️ <b>Qo'lda Yozilgan Matn</b> - Daftarga yozish uchun qulay format\n"
      "• ❓ <b>Test / Quiz / Savol-Javob</b> - Savollarga aniq javoblar va"
      " testlar\n"
      "• 🌐 <b>Language</b> - Tillar bo'yicha tarjima va grammatika\n"
      "• 📰 <b>Maqola Yozish</b> - Ilmiy va ommabop maqolalar\n"
      "• 🎨 <b>AI Rasm Yaratish</b> - Matnli tavsif bo'yicha rasm chizish\n\n"
      "👇 Pastdagi menyudan kerakli bo'limni tanlang va topshiriq mavzusini"
      " yuboring!\n\n"
      f"🎁 Sizda <b>{limit} ta bepul limit</b> bor!"
  )
  bot.send_message(
      message.chat.id,
      welcome_text,
      parse_mode="HTML",
      reply_markup=main_menu(),
  )


@bot.message_handler(func=lambda msg: True)
def handle_all_messages(message):
  chat_id = message.chat.id
  user_id = message.from_user.id
  text = message.text

  # LIMIT TEKSHIRUVI VA BO'LIMLARNI TANLASH
  sections = {
      "Slayd": ("waiting_slide", "✏️ Slayd uchun <b>mavzuni</b> kiriting:"),
      "Referat": ("waiting_referat", "📝 Referat uchun <b>mavzuni</b> kiriting:"),
      "Kurs Ishi": (
          "waiting_kurs",
          "📚 Kurs ishi uchun <b>mavzuni</b> kiriting:",
      ),
      "Mustaqil Ish": (
          "waiting_mustaqil",
          "📑 Mustaqil ish uchun <b>mavzuni</b> kiriting:",
      ),
      "Qo'lda Yozilgan": (
          "waiting_handwriting",
          "✍️ Daftarga yozish uchun moslashtiriladigan <b>mavzuni</b> kiriting:",
      ),
      "Test": (
          "waiting_quiz",
          "❓ Test, quiz yoki savolingizni kiritib yuboring:",
      ),
      "Language": (
          "waiting_lang",
          "🌐 Qaysi tildan tarjima yoki o'rganish bo'yicha yordam kerak?",
      ),
      "Maqola": ("waiting_article", "📰 Maqola uchun <b>mavzuni</b> kiriting:"),
      "Rasm": (
          "waiting_image",
          "🎨 Qanday rasm chizay? Rasm tavsifini inglizcha yoki o'zbekcha"
          " yozing:",
      ),
  }

  for key, (state_val, prompt_msg) in sections.items():
    if key in text:
      if get_user_limit(user_id) <= 0:
        bot.send_message(chat_id, "❌ Sizda bepul limitlar tugadi.")
        return
      user_states[user_id] = state_val
      bot.send_message(chat_id, prompt_msg, parse_mode="HTML")
      return

  if "Limitlarim" in text:
    limit = get_user_limit(user_id)
    bot.send_message(
        chat_id,
        f"📊 Sizning balansingizda: <b>{limit} ta</b> limit mavjud.",
        parse_mode="HTML",
    )
    return

  # PROCESS STATES
  state = user_states.get(user_id)

  if not state:
    bot.send_message(
        chat_id,
        "Iltimos, avval menyudan kerakli bo'limni tanlang!",
        reply_markup=main_menu(),
    )
    return

  user_states[user_id] = None

  # 🎨 AI RASM YARATISH
  if state == "waiting_image":
    bot.send_message(chat_id, "🎨 AI rasm chizmoqda, kuting...")
    try:
      prompt_encoded = urllib.parse.quote(text)
      image_url = f"https://pollinations.ai/p/{prompt_encoded}?width=1024&height=1024&seed=42"
      update_user_limit(user_id, -1)
      bot.send_photo(
          chat_id,
          image_url,
          caption=f"🎨 <b>Rasm tayyor!</b>\n<i>Mavzu: {text}</i>",
          parse_mode="HTML",
      )
    except Exception as e:
      bot.send_message(chat_id, f"⚠️ Rasm yaratishda xatolik: {e}")
    return

  # 📊 SLAYD YARATISH
  elif state == "waiting_slide":
    bot.send_message(chat_id, "🔄 AI slayd shakllantirmoqda...")
    prompt = f"'{text}' mavzusida 4 ta slayddan iborat taqdimot tayyorla. JSON formatda ber: [{{\"title\": \"1-slayd\", \"content\": [\"1-fakt\", \"2-fakt\"]}}]"
    try:
      response = model.generate_content(prompt)
      clean_text = response.text.strip()
      if "```json" in clean_text:
        clean_text = clean_text.split("```json")[1].split("```")[0].strip()
      elif "```" in clean_text:
        clean_text = clean_text.split("```")[1].split("```")[0].strip()

      slides_data = json.loads(clean_text)
      prs = Presentation()

      for slide_info in slides_data:
        blank = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank)

        tx = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.6), Inches(8.4), Inches(1)
        )
        p = tx.text_frame.paragraphs[0]
        p.text = slide_info.get("title", "")
        p.font.size = Pt(26)
        p.font.bold = True

        tx2 = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.8), Inches(8.4), Inches(4.5)
        )
        tf2 = tx2.text_frame
        for point in slide_info.get("content", []):
          p2 = tf2.add_paragraph()
          p2.text = f"• {point}"
          p2.font.size = Pt(18)

      filename = f"slayd_{user_id}.pptx"
      prs.save(filename)
      update_user_limit(user_id, -1)

      with open(filename, "rb") as doc:
        bot.send_document(
            chat_id,
            doc,
            caption=f"✅ <b>'{text}'</b> slaydi tayyor!",
            parse_mode="HTML",
        )
    except Exception as e:
      bot.send_message(chat_id, f"⚠️ Xatolik yuz berdi: {e}")
    return

  # 📝 MATNLI AI GENERASIYALAR (Referat, Kurs ishi, va b.)
  prompts = {
      "waiting_referat": (
          f"'{text}' mavzusida Kirish, Rejalar va Xulosadan iborat batafsil"
          " referat yoz."
      ),
      "waiting_kurs": (
          f"'{text}' mavzusida Kirish, 3 ta Asosiy bob, Amaliy qism va Xulosa"
          " mavjud bo'lgan akademik Kurs ishi rejasini va batafsil matnini"
          " yozib ber."
      ),
      "waiting_mustaqil": (
          f"'{text}' mavzusida talabalar uchun ixcham va tushunarli Mustaqil"
          " ish tayyorlab ber."
      ),
      "waiting_handwriting": (
          f"'{text}' mavzusida qo'lda daftarga ko'chirishga juda mos,"
          " soddalashtirilgan, muhim joylari ajratilgan matn tuzib ber."
      ),
      "waiting_quiz": (
          f"'{text}' mavzusida 5 ta savol-javobli test va ularning to'g'ri"
          " javoblariga batafsil izoh yoz."
      ),
      "waiting_lang": (
          f"Quyidagi matn yoki topshiriqni tahlil qilib, til o'rganish bo'yicha"
          f" batafsil yordam va tarjimasini ber: '{text}'"
      ),
      "waiting_article": (
          f"'{text}' mavzusida ilmiy-ommabop, strukturaga ega professional"
          " maqola yozib ber."
      ),
  }

  ai_prompt = prompts.get(state)
  if ai_prompt:
    bot.send_message(chat_id, "⚡️ AI ma'lumotlarni tayyorlamoqda...")
    try:
      response = model.generate_content(ai_prompt)
      update_user_limit(user_id, -1)
      bot.send_message(
          chat_id,
          f"📄 <b>Natija ({text}):</b>\n\n{response.text[:3800]}",
          parse_mode="HTML",
      )
    except Exception as e:
      bot.send_message(chat_id, f"⚠️ Xatolik: {e}")


try:
  print("--> MUKAMMAL BOT ISHGA TUSHDI <--", flush=True)
  bot.polling(non_stop=True)
except Exception as e:
  print(f"XATOLIK: {e}", flush=True)
